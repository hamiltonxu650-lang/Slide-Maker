from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


MAX_SLIDE_INCHES = 55.95

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Clean default slides
        self.prs.slides._sldIdLst.clear()
        self.dpi = 96
        self.coordinate_scale = 1.0

    def set_slide_dimensions(self, width_px, height_px, dpi=96):
        """
        Set presentation size from source pixels.

        PowerPoint has a hard limit of 56 inches per side. When the source page
        exceeds that limit, scale the whole canvas down proportionally and apply
        the same factor to every text box coordinate.
        """
        width_in = float(width_px) / float(dpi)
        height_in = float(height_px) / float(dpi)
        longest_edge = max(width_in, height_in, 0.01)
        self.coordinate_scale = min(1.0, MAX_SLIDE_INCHES / longest_edge)
        self.prs.slide_width = Inches(width_in * self.coordinate_scale)
        self.prs.slide_height = Inches(height_in * self.coordinate_scale)
        self.dpi = dpi
        return self.coordinate_scale

    def add_slide(self, background_image_path, text_data, dpi=96):
        """
        background_image_path: Path to the clean (inpainted) background image
        text_data: List of dicts with 'box', 'text', 'color', 'font_size'
        """
        # Blank slide layout is index 6 usually
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # 1. Add background image
        if background_image_path:
            slide.shapes.add_picture(
                background_image_path, 
                0, 0, 
                width=self.prs.slide_width, 
                height=self.prs.slide_height
            )

        # 2. Add text boxes at the appropriate locations
        for td in text_data:
            box = td['box']
            text = td['text']
            color = td['color']
            font_size = td['font_size']
            box_scale = float(td.get('pptx_box_scale', 1.0))
            
            # PaddleOCR box coords: [[x1, y1], [x2, y2], [x3, y3], [x4, y4]]
            x_min = min(p[0] for p in box)
            y_min = min(p[1] for p in box)
            box_width = max(p[0] for p in box) - x_min
            box_height = max(p[1] for p in box) - y_min
            scaled_width = max(box_width * box_scale, box_width)
            scaled_height = max(box_height * box_scale, box_height)
            canvas_scale = float(getattr(self, "coordinate_scale", 1.0))

            # Convert to inches
            left = Inches((x_min / dpi) * canvas_scale)
            top = Inches((y_min / dpi) * canvas_scale)
            width = Inches((scaled_width / dpi) * canvas_scale)
            height = Inches((scaled_height / dpi) * canvas_scale)

            # Create text box
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = False
            
            # Write text 
            p = tf.paragraphs[0]
            p.text = text
            
            # Format font
            p.font.size = Pt(font_size)
            # Handle potential transparency or out-of-range colors
            r = min(max(color[0], 0), 255)
            g = min(max(color[1], 0), 255)
            b = min(max(color[2], 0), 255)
            p.font.color.rgb = RGBColor(r, g, b)
            
            # Usually text is aligned left in OCR bounding boxes by default
            p.alignment = PP_ALIGN.LEFT

    def save(self, filepath):
        self.prs.save(filepath)
