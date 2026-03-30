import os
import sys
from pptx import Presentation
from PIL import Image
import io

def extract_slides_as_images(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    prs = Presentation(pptx_path)
    
    for i, slide in enumerate(prs.slides):
        print(f"Processing slide {i+1}...")
        # Method 1: Look for the largest picture on the slide
        # This is common for 'image-only' slides
        found_img = False
        largest_area = 0
        best_image_bytes = None
        best_image_ext = None
        
        for shape in slide.shapes:
            if shape.shape_type == 13: # 13 is Picture
                area = shape.width * shape.height
                if area > largest_area:
                    largest_area = area
                    best_image_bytes = shape.image.blob
                    best_image_ext = shape.image.ext
                    
        if best_image_bytes is not None:
            output_path = os.path.join(output_dir, f"slide_{i+1:03d}.{best_image_ext}")
            with open(output_path, "wb") as f:
                f.write(best_image_bytes)
            print(f"  Extracted image to {output_path}")
            found_img = True
        
        if not found_img:
            # Fallback: if no picture shape, maybe it's a background?
            # Or we might need a more advanced renderer.
            # For now, let's see if this works for the user's file.
            print(f"  Warning: No picture found on slide {i+1}")

if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "test/original.pptx"
    out_folder = "test/extracted_slides"
    extract_slides_as_images(pptx_file, out_folder)
